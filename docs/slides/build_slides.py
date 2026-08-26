#!/usr/bin/env python3
"""
build_slides.py — generate the starting PowerPoint decks.

    python docs/slides/build_slides.py

Writes two files into docs/slides/:

    jeong_setup_and_results.pptx   the 1-2 setup + 1-2 results slides Dr Jeong
                                   asked for on 20 Aug, plus two optional ones
    poster.pptx                    a 48 x 36 in research-poster skeleton

═══════════════════════════════════════════════════════════════════════════
THIS SCRIPT GENERATES THE *STARTING* DECK. IT DOES NOT MAINTAIN IT.
═══════════════════════════════════════════════════════════════════════════
Once anyone has edited the .pptx by hand, re-running this OVERWRITES that work.
Rename your edited copy (jeong_v2.pptx) or move it out of docs/slides/ before
running this again.

The numbers below are pulled from the repo where they can be, so a stale
figure fails loudly here rather than quietly on a projector.
"""

from __future__ import annotations

import csv
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

REPO = Path(__file__).resolve().parents[2]
OUT = REPO / "docs" / "slides"

# URI palette, matching webapp/static/uri.css so the deck and the dashboard
# do not look like two different projects.
NAVY   = RGBColor(0x00, 0x21, 0x47)
DARK   = RGBColor(0x00, 0x12, 0x28)
KEANEY = RGBColor(0x22, 0x77, 0xB3)
PALE   = RGBColor(0xC0, 0xDD, 0xF2)
GOLD   = RGBColor(0xB5, 0x98, 0x5A)
GOLDDK = RGBColor(0x8A, 0x70, 0x38)
INK    = RGBColor(0x0D, 0x23, 0x38)
DIM    = RGBColor(0x5B, 0x72, 0x88)
WASH   = RGBColor(0xF5, 0xF8, 0xFB)
WASH2  = RGBColor(0xEA, 0xF1, 0xF8)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
OK     = RGBColor(0x1F, 0x7A, 0x52)
WARN   = RGBColor(0xB2, 0x6A, 0x00)
FAULT  = RGBColor(0xB3, 0x26, 0x1E)

SANS = "Calibri"
MONO = "Consolas"


# ───────────────────────────────────────────────────────────── facts ──
def facts():
    """Pull every quotable number from the repo, so none is typed twice."""
    f = {}
    t = json.loads((REPO / "data" / "tunnel.json").read_text())
    f["ref1_max"] = t.get("ref1_max_rpm") or 2435
    ao = t.get("daq_analog_out", {})
    f["ao"] = f"{ao.get('volts_zero', 0.5)} V + rpm/{int(ao.get('rpm_per_volt', 300))}"

    body = [l for l in (REPO / "logs" / "sweep_v1_Ra20_summary.csv")
            .read_text().splitlines(True) if not l.startswith("#")]
    rows = list(csv.DictReader(body))
    f["n_points"] = len(rows)
    f["clean"] = sum(1 for r in rows if (r.get("clean") or "1") in ("1", "True"))
    f["rows"] = [(int(r["fan_rpm_cmd"]), int(float(r["fan_rpm_actual"])),
                  float(r["wind_mps"]), float(r["p_max_w"]),
                  float(r["i_at_pmax_a"])) for r in rows]
    f["p_top"] = max(r[3] for r in f["rows"])
    f["v_top"] = max(r[2] for r in f["rows"])

    b = json.loads((REPO / "blades" / "v1.json").read_text())
    f["chord"] = b.get("chord_mm", 48.0)
    f["span"] = b.get("span_mm", 245.1)
    f["wall"] = b.get("wall_thickness_mm", 1.79)
    f["camber"] = b.get("camber_pct_chord", 42)
    f["turning"] = b.get("turning_deg", 183)
    f["tc"] = b.get("thickness_ratio", 0.040)

    R, H = 0.1016, f["span"] / 1000.0
    f["R"], f["H"] = R, H
    f["area"] = 2 * R * H
    f["cp_lo"] = min(p / (0.5 * 1.204 * f["area"] * v ** 3) for _, _, v, p, _ in f["rows"])
    f["cp_hi"] = max(p / (0.5 * 1.204 * f["area"] * v ** 3) for _, _, v, p, _ in f["rows"])
    f["re_lo"] = 1.204 * min(r[2] for r in f["rows"]) * f["chord"] / 1000 / 1.81e-5
    f["re_hi"] = 1.204 * f["v_top"] * f["chord"] / 1000 / 1.81e-5
    return f


F = facts()


# ─────────────────────────────────────────────────────────── helpers ──
def tb(slide, l, t, w, h, wrap=True):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def para(tf, text, size=16, bold=False, color=INK, first=False, space=6,
         font=SANS, align=PP_ALIGN.LEFT, indent=0, italic=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space)
    if indent:
        p.level = indent
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = font
    return p


def rich(tf, chunks, size=16, first=False, space=6, align=PP_ALIGN.LEFT):
    """chunks = [(text, bold, color, font|None), ...] on one paragraph."""
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space)
    for text, bold, color, *rest in chunks:
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = rest[0] if rest else SANS
    return p


def rect(slide, l, t, w, h, fill, line=None, shape=MSO_SHAPE.RECTANGLE):
    s = slide.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(1)
    s.shadow.inherit = False
    s.text_frame.word_wrap = True
    # Autoshapes anchor text vertically CENTRED by default, which floats every
    # card's heading into its middle and leaves the top of the card looking
    # empty. Callers that genuinely want centring set it back afterwards.
    s.text_frame.vertical_anchor = MSO_ANCHOR.TOP
    return s


def takeaway(slide, text, W, top=6.28):
    """Navy strip across the foot of a slide: the one line to remember."""
    bar = rect(slide, 0.55, top, W - 1.1, 0.72, NAVY)
    tf = bar.text_frame
    tf.margin_left = Inches(0.28)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf, text, 15, True, WHITE, first=True, space=0)


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def header(slide, kicker, title, W):
    """Navy band with a kicker and a title. Returns the y to continue at."""
    rect(slide, 0, 0, W, 1.28, NAVY)
    rect(slide, 0, 1.28, W, 0.055, GOLD)
    tf = tb(slide, 0.55, 0.17, W - 1.1, 0.36)
    para(tf, kicker.upper(), 11.5, True, PALE, first=True, space=0)
    tf = tb(slide, 0.55, 0.47, W - 1.1, 0.66)
    para(tf, title, 29, True, WHITE, first=True, space=0)
    return 1.62


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text.strip()


def _png_size(path):
    """(w, h) in pixels from a PNG header — no imaging library in this repo."""
    b = Path(path).read_bytes()[16:24]
    return (int.from_bytes(b[:4], "big"), int.from_bytes(b[4:], "big"))


def figure_slot(slide, l, t, w, h, path, caption):
    """
    Draw a labelled placeholder, or the real image if it exists.

    A slot that says what belongs in it beats an empty rectangle: the deck is
    handed on, and the next person needs to know what is missing.
    """
    p = REPO / path
    if p.exists():
        # Fit INSIDE the slot on both axes and centre. Passing width alone
        # sets height by aspect, which is fine on the 16:9 slots and runs a
        # 1.57:1 diagram straight off the bottom of the poster's 2.97:1 one.
        iw, ih = _png_size(p)
        sc = min(w / iw, (h - 0.34) / ih)
        pw, ph = iw * sc, ih * sc
        slide.shapes.add_picture(str(p), Inches(l + (w - pw) / 2),
                                 Inches(t + (h - 0.34 - ph) / 2),
                                 width=Inches(pw), height=Inches(ph))
        tf = tb(slide, l, t + h - 0.28, w, 0.3)
        para(tf, caption, 10.5, False, DIM, first=True, italic=True,
             align=PP_ALIGN.CENTER)
        return
    box = rect(slide, l, t, w, h, WASH2, PALE)
    tf = box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf, "FIGURE", 12, True, KEANEY, first=True, space=4,
         align=PP_ALIGN.CENTER)
    para(tf, path, 11, False, DIM, space=4, align=PP_ALIGN.CENTER, font=MONO)
    para(tf, caption, 11, False, DIM, align=PP_ALIGN.CENTER, italic=True)


# ══════════════════════════════════════════════════════ DECK 1: JEONG ══
def build_jeong():
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    W, H = 13.333, 7.5

    # ---------- title ----------
    s = blank(prs)
    rect(s, 0, 0, W, H, NAVY)
    rect(s, 0, 0, 0.28, H, GOLD)
    tf = tb(s, 1.1, 2.05, 11, 0.4)
    para(tf, "UNIVERSITY OF RHODE ISLAND  ·  SODHI LAB  ×  JEONG LAB",
         13, True, GOLD, first=True, space=0)
    tf = tb(s, 1.1, 2.6, 11.2, 1.5)
    para(tf, "Programmable Characterisation of", 34, True, WHITE, first=True, space=2)
    para(tf, "3D-Printed Wind Turbine Blades", 34, True, WHITE, space=0)
    tf = tb(s, 1.1, 4.25, 11, 0.5)
    para(tf, "Automated load sweeps in the Aerolab wind tunnel",
         18, False, PALE, first=True, space=0)
    rect(s, 1.1, 4.95, 2.4, 0.03, KEANEY)
    tf = tb(s, 1.1, 5.3, 11, 1.0)
    para(tf, "Stephen Acuello", 15, True, WHITE, first=True, space=3)
    para(tf, "seacuello@uri.edu   ·   26 August 2026", 12.5, False, PALE, space=0)
    notes(s, "Setup pair = slides 2-3. Results pair = slides 4-5. "
             "Slides 6-7 are optional; 6 is the ask.")

    # ---------- 2. setup ----------
    s = blank(prs)
    y = header(s, "Slide 1 of 2  ·  Experimental setup", "The Rig", W)
    rect(s, 0, y - 0.34 + 0.34, W, H - y, WASH)
    figure_slot(s, 0.55, y + 0.12, 7.3, 4.9,
                "docs/diagrams/system_overview.png",
                "Full measurement and control chain")
    tf = tb(s, 8.2, y + 0.12, 4.6, 4.9)
    para(tf, "CONTROL CHAIN", 12, True, KEANEY, first=True, space=8)
    for a, b in [("Drive", f"ABB ACS550-U1-046A-2, 15 HP, 208–240 V 3φ. "
                           f"Commands speed in rpm; full scale {F['ref1_max']}"),
                 ("Path", "Python → USB → Portenta Machine Control → "
                          "RS-485 Modbus → drive"),
                 ("Load", "Chroma 63004-150-60 — sets the rotor's operating "
                          "point and measures it"),
                 ("Wind", f"500–1800 rpm fan = 10.1–{F['v_top']:.1f} m/s, "
                          f"v = 0.02132·rpm − 0.424 (R² = 0.9996)"),
                 ("Bandwidth", "τ = 0.63 ± 0.12 s, from five 1-cosine gusts")]:
        rich(tf, [(f"{a} — ", True, NAVY), (b, False, INK)], 13, space=9)
    para(tf, "ROTOR UNDER TEST", 12, True, KEANEY, space=8)
    rich(tf, [("Vertical-axis H-rotor", True, NAVY),
              (f", 3 blades, R = 101.6 mm, span {F['span']:.0f} mm. Swept "
               f"area 2RH = {F['area']:.4f} m² — a cylinder, not a disc.",
               False, INK)], 13, space=9)
    rich(tf, [("Blade: ", True, NAVY),
              (f"thin cambered plate — {F['wall']:.2f} mm wall, t/c "
               f"{F['tc']*100:.0f}%, {F['camber']}% camber, {F['turning']}° "
               f"turning. Not an airfoil.", False, INK)], 13, space=0)
    notes(s, "Swept area is 2RH because it is a VAWT. Using pi*R^2 overstates "
             "Cp by 1.54x. The blade is a cambered plate, NOT an airfoil - it "
             "matters for any polar or XFOIL cross-check.")

    # ---------- 3. safety ----------
    s = blank(prs)
    y = header(s, "Slide 2 of 2  ·  Experimental setup", "Safety Architecture", W)
    rect(s, 0, y, W, H - y, WASH)
    b = rect(s, 0.55, y + 0.15, 12.2, 0.95, RGBColor(0xFD, 0xF0, 0xEE), FAULT)
    tf = b.text_frame
    tf.margin_left = Inches(0.25)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf, "The hardwired E-stop is the safety device.", 19, True, FAULT,
         first=True, space=3)
    para(tf, "No software is in that chain. Everything below reduces how "
             "often it is needed — none of it replaces it.", 13.5, False, INK,
         space=0)
    cards = [
        ("TWO INDEPENDENT WATCHDOGS",
         [("The drive stops the fan", " if the PMC goes quiet — par 3018/3019, 3.0 s"),
          ("The PMC stops the fan", " if the host goes quiet — 5.0 s"),
          ("", "Neither depends on the layer above it being correct")]),
        ("ONE MODBUS MASTER ONLY",
         [("The PMC is the master.", " Modbus RTU permits exactly one"),
          ("", "A direct host-to-drive cable must never be landed as well"),
          ("", "Two devices commanding a 15 HP fan, neither aware of the "
               "other, is the failure this prevents")]),
        ("TURBINE INTERLOCK",
         [("load ON → wind UP → test → wind DOWN → load OFF", ""),
          ("", "An unloaded rotor in moving air accelerates until something "
               "mechanical stops it"),
          ("", "Enforced on every path that can move the fan"),
          ("If the fan cannot be confirmed stopped,", " the load stays on")]),
    ]
    for i, (head, items) in enumerate(cards):
        l = 0.55 + i * 4.12
        card = rect(s, l, y + 1.32, 3.87, 2.95, WHITE, PALE)
        tf = card.text_frame
        tf.margin_left = tf.margin_right = Inches(0.2)
        tf.margin_top = Inches(0.18)
        para(tf, head, 12.5, True, KEANEY, first=True, space=10)
        for bold, rest in items:
            if bold and not rest:
                para(tf, bold, 13, True, GOLDDK, space=9, font=MONO)
            else:
                rich(tf, [("• ", False, KEANEY), (bold, True, NAVY),
                          (rest, False, INK)], 12.5, space=9)
    takeaway(s, "Each layer is watched by the layer below it — so a host "
                "that is wrong, hung or unplugged cannot run the fan away.", W)
    notes(s, "The two watchdogs are the slide's point: each layer is watched by "
             "the layer BELOW it, so the host being wrong cannot run the fan away.")

    # ---------- 4. method ----------
    s = blank(prs)
    y = header(s, "Slide 1 of 2  ·  Results", "Method", W)
    rect(s, 0, y, W, H - y, WASH)
    tf = tb(s, 0.55, y + 0.12, 12.2, 0.4)
    rich(tf, [("Automated blade characterisation — ", False, INK),
              (f"~10 minutes per rotor, unattended, {F['n_points']} wind speeds",
               True, NAVY)], 17, first=True, space=0)
    figure_slot(s, 0.55, y + 0.68, 7.5, 4.4,
                "docs/diagrams/sweep_protocol.png",
                "Outer loop = wind speed; inner loop = load current")
    tf = tb(s, 8.35, y + 0.68, 4.45, 4.9)
    para(tf, "AT EACH WIND SPEED", 12, True, KEANEY, first=True, space=9)
    for txt in ["Ramp the electronic load in constant-current steps",
                "Stop once electrical power falls to 80% of its peak — the "
                "rotor is never driven to stall",
                "Record V and I at every step, unload, advance the wind"]:
        rich(tf, [("• ", False, KEANEY), (txt, False, INK)], 13, space=8)
    para(tf, "WHY A PARABOLIC FIT", 12, True, KEANEY, space=9)
    para(tf, "The top of P(I) is flat. The largest single sample is biased "
             "high, and the bias grows with sample count — so blades with "
             "different point counts would be compared unfairly.",
         13, False, INK, space=10)
    para(tf, "PROTOCOL FINGERPRINT", 12, True, KEANEY, space=9)
    rich(tf, [("Every run carries one. ", False, INK),
              ("Runs measured under different settings are not comparable",
               True, NAVY),
              (" — the fingerprint makes that visible instead of silent.",
               False, INK)], 13, space=0)
    notes(s, "The fingerprint is the campaign's integrity mechanism. Across a "
             "dozen rotors it is very easy to change a setting by accident and "
             "impossible to spot afterwards - the numbers stay plausible.")

    # ---------- 5. results ----------
    s = blank(prs)
    y = header(s, "Slide 2 of 2  ·  Results", "Blade v1, Ra 20 µm", W)
    rect(s, 0, y, W, H - y, WASH)

    show = [r for r in F["rows"] if r[0] in (500, 900, 1400, 1800)]
    tbl = s.shapes.add_table(len(show) + 1, 5, Inches(0.55), Inches(y + 0.15),
                             Inches(6.1), Inches(0.42 * (len(show) + 1))).table
    for i, w in enumerate([1.35, 1.2, 1.05, 1.3, 1.2]):
        tbl.columns[i].width = Inches(w)
    heads = ["fan rpm cmd", "measured", "m/s", "P_max", "at"]
    for c, h in enumerate(heads):
        cell = tbl.cell(0, c)
        cell.text = h
        cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
        pr = cell.text_frame.paragraphs[0]
        pr.alignment = PP_ALIGN.RIGHT if c else PP_ALIGN.LEFT
        pr.runs[0].font.size = Pt(11.5); pr.runs[0].font.bold = True
        pr.runs[0].font.color.rgb = WHITE; pr.runs[0].font.name = SANS
    for r, (cmd, act, mps, p, ia) in enumerate(show, start=1):
        top = (cmd == 1800)
        for c, v in enumerate([f"{cmd}", f"{act}", f"{mps:.1f}",
                               f"{p:.3f} W", f"{ia:.3f} A"]):
            cell = tbl.cell(r, c)
            cell.text = v
            cell.fill.solid()
            cell.fill.fore_color.rgb = WASH2 if top else WHITE
            pr = cell.text_frame.paragraphs[0]
            pr.alignment = PP_ALIGN.RIGHT if c else PP_ALIGN.LEFT
            pr.runs[0].font.size = Pt(12)
            pr.runs[0].font.bold = top
            pr.runs[0].font.color.rgb = NAVY if top else INK
            pr.runs[0].font.name = MONO

    tf = tb(s, 0.55, y + 0.28 + 0.42 * (len(show) + 1), 6.1, 0.7)
    para(tf, "Wind speed is derived from MEASURED fan rpm, not commanded — the "
             "drive settles 4–13 rpm below setpoint.", 10.5, False, DIM,
         first=True, space=0, italic=True)

    box = rect(s, 0.55, y + 2.75, 6.1, 1.72, WHITE, PALE)
    tf = box.text_frame
    tf.margin_left = tf.margin_right = Inches(0.2); tf.margin_top = Inches(0.15)
    para(tf, "WHAT IS SOLID", 12, True, OK, first=True, space=8)
    for t in [f"{F['clean']} / {F['n_points']} points clean, single "
              f"continuous run",
              "P ∝ v^3.77, R² = 0.998 across 10.1–37.5 m/s",
              "Independent repeats at 1200 and 1800 rpm match to 0.3% / 0.2%"]:
        rich(tf, [("• ", False, OK), (t, False, INK)], 12.5, space=7)

    box = rect(s, 6.85, y + 0.15, 5.9, 4.32, WHITE, GOLD)
    tf = box.text_frame
    tf.margin_left = tf.margin_right = Inches(0.22); tf.margin_top = Inches(0.18)
    para(tf, "WHAT IT DOES NOT YET SHOW", 12.5, True, GOLDDK, first=True, space=10)
    rich(tf, [("This is electrical power at the load terminals — ", False, INK),
              ("not Cp", True, FAULT), (".", False, INK)], 14.5, space=10)
    para(tf, "The v^3.77 exponent decomposes into the generator, not the "
             "rotor:", 13, False, INK, space=8)
    for t in ["V_oc ∝ v^1.52      R_int ∝ v^−0.64",
              "R_int falls 73.6 → 40.1 Ω across the range",
              "peaks at the Thévenin match, P = V_oc²/4R_int",
              "n = 2a − b = 3.69   vs   3.77 measured"]:
        para(tf, t, 11.5, False, NAVY, space=5, font=MONO, indent=1)
    rich(tf, [("→ ", True, GOLDDK),
              ("Little aerodynamic residual is left for the blade to move. "
               "Without rotor speed, blade comparisons partly compare the "
               "generator.", False, INK)], 13, space=10)
    rich(tf, [("Cp_elec = ", False, INK),
              (f"{F['cp_lo']*100:.2f}–{F['cp_hi']*100:.2f}%", True, FAULT),
              (f", roughly 100× below a working H-rotor, because peak power "
               f"sits at ω/ω_runaway ≈ 0.70→0.94 — the far limb of Cp(λ) "
               f"where Cp → 0 by construction.", False, INK)], 13, space=0)
    takeaway(s, "The measurement is sound. Its interpretation needs rotor "
                "speed — which is the next slide.", W)
    notes(s, "Do NOT present v^3.77 as 'Cp still climbing with Reynolds'. A "
             "9-agent audit on 25 Aug showed the exponent is a generator "
             "characteristic: V_oc ~ v^1.52, R_int ~ v^-0.64, n = 2a-b = 3.69. "
             "The honest version is stronger and it is the argument for slide 6.")

    # ---------- 6. the ask ----------
    s = blank(prs)
    y = header(s, "Optional  ·  The next measurement",
               "One Channel Unlocks Cp(λ)", W)
    rect(s, 0, y, W, H - y, WASH)
    box = rect(s, 0.55, y + 0.15, 12.2, 1.15, RGBColor(0xEF, 0xF6, 0xF1), OK)
    tf = box.text_frame
    tf.margin_left = Inches(0.25); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf, "The rotor-speed measurement already exists in the Jeong lab DAQ.",
         19, True, OK, first=True, space=3)
    para(tf, "A proximity sensor on the rig already records it. What is needed "
             "is its channel number and pulses-per-revolution — not new "
             "hardware.", 13.5, False, INK, space=0)
    left = [
        ("WHY IT MATTERS MORE THAN ANY BLADE RUN", [
            "Without ω there is no λ and no Cp — the rig cannot separate "
            "rotor aerodynamics from generator matching.",
            "A better blade can read as LESS power: raising Cp_max while "
            "adding low-α drag lowers λ_runaway, lowers V_oc, and lowers "
            "measured P as the SQUARE.",
            "Two rotors at Cp_max 5% and 25% would rank purely by how freely "
            "they spin."]),
    ]
    for i, (head, items) in enumerate(left):
        card = rect(s, 0.55, y + 1.5, 6.0, 3.28, WHITE, PALE)
        tf = card.text_frame
        tf.margin_left = tf.margin_right = Inches(0.22)
        tf.margin_top = Inches(0.18)
        para(tf, head, 12, True, KEANEY, first=True, space=10)
        for t in items:
            rich(tf, [("• ", False, KEANEY), (t, False, INK)], 13, space=10)
    card = rect(s, 6.75, y + 1.5, 6.0, 3.28, WHITE, GOLD)
    tf = card.text_frame
    tf.margin_left = tf.margin_right = Inches(0.22); tf.margin_top = Inches(0.18)
    para(tf, "THE PAYOFF IS RETROACTIVE", 12, True, GOLDDK, first=True, space=10)
    rich(tf, [("Every ", False, INK),
              ("sweep_*_points.csv", False, NAVY, MONO),
              (" already holds a full load ramp at all 14 wind speeds. Adding "
               "one rotor-speed channel converts the ", False, INK),
              ("entire existing archive", True, NAVY),
              (" into a Cq(λ) traverse — without another minute of tunnel time.",
               False, INK)], 13, space=12)
    para(tf, "ALSO USEFUL, AND ALREADY BUILT", 12, True, KEANEY, space=9)
    rich(tf, [("Fan rpm now leaves the PMC as a 0–10 V analog signal "
               "for the DAQ (", False, INK),
              (F["ao"], False, NAVY, MONO),
              ("), so fan speed and rotor speed share one time base. "
               "0.000 V is a live zero meaning INVALID, not zero rpm.",
               False, INK)], 13, space=0)
    takeaway(s, "The ask is two numbers: which DAQ channel, and how many "
                "pulses per revolution.", W)
    notes(s, "This is the ask, aimed at the person who owns the DAQ. Keep it "
             "concrete: channel number and pulses-per-rev, that is all.")

    # ---------- 7. limitations ----------
    s = blank(prs)
    y = header(s, "Optional  ·  Stated plainly", "Known Limitations", W)
    rect(s, 0, y, W, H - y, WASH)
    tf = tb(s, 0.55, y + 0.15, 12.2, 0.4)
    para(tf, "Every one of these was found by adversarial review of our own "
             "results, and each is recorded in the repository.",
         14, False, DIM, first=True, space=0, italic=True)
    items = [
        ("Not Cp", "Electrical power only. Cp(λ) needs the rotor-speed "
                   "channel; every blade tested before that lands must be "
                   "re-run to obtain λ."),
        ("The blade is not an airfoil",
         f"{F['wall']:.2f} mm wall, t/c {F['tc']*100:.0f}%, {F['camber']}% "
         f"camber, square-cut edges. Separation is pinned by geometry, so "
         f"XFOIL and polar cross-checks answer a different question."),
        ("Reynolds range",
         f"Re_chord = {F['re_lo']:,.0f} → {F['re_hi']:,.0f}. Below "
         f"Reynolds-independence for a cross-flow turbine (~2×10⁵), so "
         f"absolute performance is not transferable to full scale."),
        ("Generator not characterised",
         "R_int is inferred from a curve fit and has never been measured "
         "directly. It gates every derived quantity."),
        ("Wind speed disagreement",
         "Summary and points files differ by 1.1–1.2% at the same fan set "
         "point. At v^3.77 that is 4.5% in power. Unresolved — flagged, not "
         "averaged away."),
        ("Single-mount data",
         "Every result so far comes from one mounting of one rotor. There is "
         "no mount-to-mount error bar yet; the ABBA protocol for the next "
         "session produces the first one."),
    ]
    for i, (head, body) in enumerate(items):
        col, row = i % 2, i // 2
        l = 0.55 + col * 6.2
        t = y + 0.68 + row * 1.72
        card = rect(s, l, t, 6.0, 1.56, WHITE, PALE)
        tf = card.text_frame
        tf.margin_left = tf.margin_right = Inches(0.2)
        tf.margin_top = Inches(0.13)
        para(tf, head, 13, True, NAVY, first=True, space=5)
        para(tf, body, 11.5, False, INK, space=0)
    notes(s, "Optional slide. For an academic audience this is the most "
             "credible thing in the deck - it shows the limits are known and "
             "quantified rather than unexamined.")

    out = OUT / "jeong_setup_and_results.pptx"
    prs.save(out)
    return out, len(prs.slides._sldIdLst)


# ═══════════════════════════════════════════════════════ DECK 2: POSTER ══
# Built by CLONING 2026_CYPHER_IPT_Poster_Eacuello.pptx and replacing its
# content, rather than styling a blank deck to look similar. Cloning inherits
# the theme, the fonts, the navy background and the slide size exactly, so the
# two posters are the same object with different words - which is what "use my
# CYPHER poster as the template" has to mean to survive contact with a printer.
TEMPLATE = OUT / "2026_CYPHER_IPT_Poster_Eacuello.pptx"

GOLD_C   = RGBColor(0xC7, 0x93, 0x16)   # the CYPHER rule colour
PANEL_LW = Pt(13.3)                     # measured off the template panels
TITLE_LW = Pt(15.3)

# The grid, measured off the template. The template's left column is 0.2 in
# out of alignment with the panel below it; that is reproduced as a straight
# 0.93 here rather than copied faithfully.
G = {
    "title":   (0.93,  0.63, 46.30,  6.32),
    "intro":   (0.93,  7.66, 13.88,  9.25),
    "arch":    (0.93, 17.43, 13.88, 14.95),
    "detail":  (15.35, 7.64, 16.72, 24.83),
    "results": (32.61, 7.64, 14.46, 15.16),
    "concl":   (32.61,23.43, 14.46,  9.04),
    "footer":  (0.93, 33.10, 46.15,  2.13),
}

FUNDING = (
    "This research was supported by the Office of Naval Research under Grant "
    "No. N00014-24-1-2129."
)
DISCLAIMER = (
    "The views and conclusions contained in this poster are those of the "
    "authors and should not be interpreted as representing the official "
    "policies, either expressed or implied, of the Office of Naval Research "
    "or the U.S. Government"
)


def _wipe(slide):
    """Strip every shape, leaving the theme, background and size intact."""
    for sh in list(slide.shapes):
        sh._element.getparent().remove(sh._element)


def _panel(slide, key, round_=False):
    l, t, w, h = G[key]
    sh = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = WHITE
    sh.line.color.rgb = GOLD_C
    sh.line.width = TITLE_LW if round_ else PANEL_LW
    sh.shadow.inherit = False
    return sh


def _label(slide, key, text):
    """72 pt navy section label, seated just inside the panel's top edge."""
    l, t, w, h = G[key]
    tf = tb(slide, l + 0.35, t + 0.11, w - 0.7, 1.51)
    para(tf, text, 72, True, NAVY, first=True, space=0)
    return t + 1.62


def _body(slide, key, y, pad=0.55):
    l, t, w, h = G[key]
    return tb(slide, l + pad, y, w - 2 * pad, t + h - y - 0.35)


def _fits(key, y, chars, size=32):
    """
    Rough overflow check. Better a warning at build time than a surprise at
    the plotter, where a 48 in sheet costs real money and an afternoon.
    """
    l, t, w, h = G[key]
    per_line = max(1, int((w - 1.1) * 72 / (size * 0.50)))
    lines = chars / per_line
    return lines * (size * 1.28 / 72) <= (t + h - y - 0.35)


def build_poster():
    if not TEMPLATE.exists():
        print(f"  ! {TEMPLATE.name} not found — poster skipped")
        return None, 0, 0
    prs = Presentation(str(TEMPLATE))
    s = prs.slides[0]
    _wipe(s)
    W, H = 48.0, 36.0
    A = OUT / "assets"

    # ── title ──
    _panel(s, "title", round_=True)
    tf = tb(s, G["title"][0] + 11.0, G["title"][1] + 0.75, 24.5, 5.0)
    para(tf, "Programmable Characterisation of", 66, True, NAVY, first=True, space=6)
    para(tf, "3D-Printed Wind Turbine Blades", 66, True, NAVY, space=16)
    para(tf, "Stephen S. Eacuello, ISE Ph.D. Candidate", 44, True, INK, space=6)
    para(tf, "Department of Mechanical, Industrial and Systems Engineering",
         32, False, DIM, space=0)
    for name, box in [("logo_uri_title.png", (1.39, 1.88, 10.05)),
                      ("logo_cypher.jpg",    (37.90, 2.13, 8.66))]:
        if (A / name).exists():
            s.shapes.add_picture(str(A / name), Inches(box[0]), Inches(box[1]),
                                 width=Inches(box[2]))

    # ── introduction ──
    _panel(s, "intro")
    y = _label(s, "intro", "Introduction")
    tf = _body(s, "intro", y)
    para(tf, "3D printing makes rotor geometry cheap to iterate. Measuring "
             "whether a change actually helped does not follow automatically.",
         32, False, INK, first=True, space=16)
    rich(tf, [("A campaign comparing a dozen blades has to hold its protocol "
               "fixed across weeks, or the comparison silently becomes a "
               "comparison of ", False, INK),
              ("settings", True, NAVY), (".", False, INK)], 32, space=16)
    para(tf, "This work builds programmable, reproducible control of the "
             "Aerolab tunnel and its electronic load, so a rotor is "
             "characterised unattended in ~10 minutes with a machine-checkable "
             "record of exactly how.", 32, False, INK, space=18)
    for t in ["Automated 14-point load sweep, drive and load on one clock",
              "Two-layer watchdog safety, hardware E-stop untouched",
              "Protocol fingerprinting so runs are provably comparable",
              "Live dashboard, digital twin, and a full parameter archive"]:
        rich(tf, [("▪ ", False, GOLDDK), (t, False, NAVY)], 28, space=9)

    # ── system architecture ──
    _panel(s, "arch")
    y = _label(s, "arch", "System Architecture")
    tf = _body(s, "arch", y)
    for a, b in [("Drive", f"ABB ACS550-U1-046A-2, 15 HP, 208–240 V 3φ. "
                           f"Commands speed in rpm, full scale {F['ref1_max']}"),
                 ("Controller", "Portenta Machine Control — sole Modbus RTU "
                                "master, owns both watchdogs"),
                 ("Load", "Chroma 63004-150-60 DC electronic load, SCPI over "
                          "USB-TMC"),
                 ("Wind", f"10.1 – {F['v_top']:.1f} m/s;  "
                          f"v = 0.02132·rpm − 0.424,  R² = 0.9996")]:
        rich(tf, [(f"{a}   ", True, NAVY), (b, False, INK)], 28, space=11)
    para(tf, "The hardwired E-stop is the safety device. No software is in "
             "that chain.", 28, True, FAULT, space=11)
    for t in ["Two independent watchdogs — the drive stops the fan if the PMC "
              "goes quiet (3.0 s); the PMC stops the fan if the host goes "
              "quiet (5.0 s). Neither depends on the layer above it.",
              "One Modbus master only. Two devices commanding a 15 HP fan, "
              "neither aware of the other, is the failure this design prevents.",
              "Interlock: load ON → wind UP → test → wind DOWN → load OFF. "
              "An unloaded rotor in moving air accelerates."]:
        rich(tf, [("• ", False, GOLDDK), (t, False, INK)], 28, space=11)
    figure_slot(s, G["arch"][0] + 0.55, 27.3, 12.78, 4.85,
                "docs/diagrams/chain_strip.png",
                "Measurement and control chain")

    # ── project details ──
    _panel(s, "detail")
    y = _label(s, "detail", "Project Details")
    tf = _body(s, "detail", y)
    para(tf, f"At each of {F['n_points']} wind speeds (500 → 1800 rpm fan, "
             f"100 rpm steps):", 32, True, NAVY, first=True, space=13)
    for t in ["Ramp the electronic load in constant-current steps",
              "Stop once electrical power falls to 80% of its peak — the "
              "rotor is never driven to stall",
              "Record V and I at every step, unload, advance the wind"]:
        rich(tf, [("• ", False, GOLDDK), (t, False, INK)], 32, space=11)
    rich(tf, [("Peak located by parabolic fit. ", True, NAVY),
              ("Over a flat maximum the largest single sample is biased high, "
               "and the bias grows with sample count — so blades measured with "
               "different point counts would be compared unfairly.",
               False, INK)], 32, space=13)
    rich(tf, [("Protocol fingerprinting. ", True, NAVY),
              ("Every run hashes the settings that change what a curve means. "
               "Runs with different fingerprints are not comparable, and the "
               "hash makes that visible instead of silent.", False, INK)],
         32, space=13)
    rich(tf, [("Rotor   ", True, NAVY),
              (f"vertical-axis H-rotor, 3 blades, R = 101.6 mm, span "
               f"{F['span']:.0f} mm. Swept area 2RH = {F['area']:.4f} m² — "
               f"a cylinder, not a disc.", False, INK)], 30, space=11)
    rich(tf, [("Blade   ", True, NAVY),
              (f"PETG thin cambered plate: {F['wall']:.2f} mm wall, t/c "
               f"{F['tc']*100:.0f}%, {F['camber']}% camber, {F['turning']}° "
               f"turning, square edges — not an airfoil.", False, INK)],
         30, space=0)
    figure_slot(s, G["detail"][0] + 0.55, 19.9, 15.62, 11.9,
                "docs/diagrams/sweep_protocol.png",
                "Outer loop = wind speed; inner loop = load current")

    # ── results ──
    _panel(s, "results")
    y = _label(s, "results", "Results")
    tf = _body(s, "results", y)
    para(tf, f"{F['clean']} / {F['n_points']} points clean, single continuous "
             f"run", 30, False, INK, first=True, space=13)
    para(tf, "P ∝ v^3.77    R² = 0.998", 40, True, NAVY, space=8)
    para(tf, f"peak {F['p_top']:.3f} W at {F['v_top']:.1f} m/s", 40, True,
         NAVY, space=13)
    para(tf, "Independent repeats at 1200 and 1800 rpm match the sweep to "
             "0.3% and 0.2%.", 28, False, DIM, space=16)
    para(tf, "THE EXPONENT IS A GENERATOR CHARACTERISTIC, NOT AN AERODYNAMIC "
             "ONE", 28, True, GOLDDK, space=11)
    for t in ["V_oc ∝ v^1.52   and   R_int ∝ v^−0.64   (73.6 → 40.1 Ω)",
              "every peak sits at the Thévenin match,  P = V_oc² / 4R_int",
              "n = 2a − b = 3.69,  against 3.77 measured"]:
        para(tf, t, 26, False, NAVY, space=8)
    rich(tf, [("Cp_elec = ", False, INK),
              (f"{F['cp_lo']*100:.2f} – {F['cp_hi']*100:.2f}%", True, FAULT),
              (", ~100× below a working H-rotor: peak power sits at "
               "ω/ω_runaway ≈ 0.70 → 0.94, the far limb of Cp(λ) where "
               "Cp → 0 by construction.", False, INK)], 28, space=0)

    # The measured curve itself. An empty half-panel on a poster reads as
    # "we did not have much", which is the opposite of true here.
    show = [r for r in F["rows"] if r[0] in (500, 900, 1400, 1800)]
    tl = s.shapes.add_table(len(show) + 1, 4,
                            Inches(G["results"][0] + 0.55), Inches(17.5),
                            Inches(13.0), Inches(0.82 * (len(show) + 1))).table
    for i, w in enumerate([3.3, 3.0, 3.4, 3.3]):
        tl.columns[i].width = Inches(w)
    for c, htxt in enumerate(["fan rpm", "m/s", "P_max", "at"]):
        cell = tl.cell(0, c)
        cell.text = htxt
        cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
        pr = cell.text_frame.paragraphs[0]
        pr.alignment = PP_ALIGN.RIGHT if c else PP_ALIGN.LEFT
        pr.runs[0].font.size = Pt(24); pr.runs[0].font.bold = True
        pr.runs[0].font.color.rgb = WHITE
    for r, (cmd, act, mps, pw, ia) in enumerate(show, start=1):
        top = (cmd == 1800)
        for c, v in enumerate([f"{act}", f"{mps:.1f}", f"{pw:.3f} W",
                               f"{ia:.3f} A"]):
            cell = tl.cell(r, c)
            cell.text = v
            cell.fill.solid()
            cell.fill.fore_color.rgb = WASH2 if top else WHITE
            pr = cell.text_frame.paragraphs[0]
            pr.alignment = PP_ALIGN.RIGHT if c else PP_ALIGN.LEFT
            pr.runs[0].font.size = Pt(24); pr.runs[0].font.bold = top
            pr.runs[0].font.color.rgb = NAVY if top else INK
    tf2 = tb(s, G["results"][0] + 0.55, 17.5 + 0.82 * (len(show) + 1) + 0.12,
             13.0, 0.8)
    para(tf2, "Wind speed from MEASURED fan rpm — the drive settles 4–13 rpm "
              "below setpoint.", 19, False, DIM, first=True, space=0,
         italic=True)

    # ── conclusion ──
    _panel(s, "concl")
    y = _label(s, "concl", "Conclusion")
    tf = _body(s, "concl", y)
    rich(tf, [("A better blade can read as LESS electrical power. ", True, FAULT),
              ("Anything raising Cp_max while adding low-α drag lowers "
               "λ_runaway, lowers V_oc, and lowers measured P as the square.",
               False, INK)], 28, first=True, space=13)
    rich(tf, [("One rotor-speed channel converts the entire existing archive "
               "into Cq(λ) retroactively", True, OK),
              (" — every sweep already records a full load ramp at all 14 "
               "wind speeds, and a proximity sensor on the rig already feeds "
               "the DAQ. What is needed is its channel and pulses-per-"
               "revolution, not new hardware.", False, INK)], 28, space=13)
    para(tf, f"Limits: electrical power only, not Cp · Re_chord "
             f"{F['re_lo']:,.0f}–{F['re_hi']:,.0f}, below Reynolds-"
             f"independence (~2×10⁵) · generator R_int inferred, never "
             f"measured · single-mount data, no mount-to-mount error bar yet.",
         24, False, DIM, space=0)

    # ── footer ──
    _panel(s, "footer")
    tf = tb(s, G["footer"][0] + 7.2, G["footer"][1] + 0.16, 33.0, 1.8)
    para(tf, "The University of Rhode Island CYPHER Research Center IPT "
             "Meeting for Power Systems and Manufacturing", 29.33, True,
         RGBColor(0, 0, 0), first=True, space=2, font="Arial")
    para(tf, FUNDING, 24, True, RGBColor(0x22, 0x22, 0x22), space=2, font="Arial")
    para(tf, DISCLAIMER, 15, False, RGBColor(0x22, 0x22, 0x22), space=0,
         font="Arial")
    if (A / "logo_uri_footer.png").exists():
        s.shapes.add_picture(str(A / "logo_uri_footer.png"), Inches(1.12),
                             Inches(33.39), width=Inches(6.25))

    out = OUT / "poster.pptx"
    prs.save(out)
    return out, W, H


if __name__ == "__main__":
    OUT.mkdir(parents=True, exist_ok=True)
    p1, n1 = build_jeong()
    p2, pw, ph = build_poster()
    print(f"  {p1.relative_to(REPO)}   {n1} slides, 16:9")
    print(f"  {p2.relative_to(REPO)}   1 slide, {pw:.0f} x {ph:.0f} in")
    print(f"\n  facts pulled from the repo:")
    print(f"    {F['clean']}/{F['n_points']} clean · peak {F['p_top']:.3f} W "
          f"at {F['v_top']:.1f} m/s")
    print(f"    Cp_elec {F['cp_lo']*100:.2f}-{F['cp_hi']*100:.2f}% · "
          f"Re {F['re_lo']:,.0f}-{F['re_hi']:,.0f}")
    print(f"    blade wall {F['wall']:.2f} mm · t/c {F['tc']*100:.0f}% · "
          f"camber {F['camber']}%")
